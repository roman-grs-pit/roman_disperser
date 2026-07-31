"""Build the v0.13.0 rerun before/after addendum deck.

Addendum to the wave-1b draft deck (line_grid_wave1b_draft.pptx), in the same
visual style (helpers copied from its build_slides.py so the decks read as one
set). One slide per question of the 2026-07-31 before/after briefing; every
number is quoted from that briefing / summary.json — regenerate
before_after.py outputs first if the stores change.

Output goes to the S3 analysis dir, NOT the repo (Nikhil's 2026-07-31
decision: decks live with the data). python-pptx runs from an ephemeral env:

    pixi exec --spec python=3.12 --spec python-pptx --spec pillow \
        python build_rerun_slides.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ANALYSIS = Path("/mnt/roman-science/grs/line-tests-20260731/analysis")
OUT = ANALYSIS / "line_grid_v0130_rerun_addendum.pptx"

# --- style constants: identical to slides/build_slides.py (wave-1b deck) ---
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
BODY_TOP = Inches(1.25)
INK = RGBColor(0x20, 0x20, 0x20)
MUTED = RGBColor(0x59, 0x59, 0x59)
ACCENT = RGBColor(0x00, 0x50, 0x80)
HDR_BG = RGBColor(0xE8, 0xEE, 0xF4)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide(title=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(MARGIN, Inches(0.35), SLIDE_W - 2 * MARGIN,
                                  Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    return s


def add_text(slide, lines, left=MARGIN, top=BODY_TOP, width=None,
             size=14, mono=False, color=INK, line_gap=6):
    """lines: list of (text, level) or plain strings (level 0)."""
    width = width or (SLIDE_W - 2 * MARGIN)
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = lvl
        p.space_after = Pt(line_gap)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            if mono:
                r.font.name = "Consolas"
    return tb


def add_figure(slide, path, caption=None, max_h_in=5.35, left=None,
               max_w_in=12.2, top=None):
    img = Image.open(path)
    w_px, h_px = img.size
    scale = min(max_w_in / (w_px / 150), max_h_in / (h_px / 150))
    w_in, h_in = (w_px / 150) * scale, (h_px / 150) * scale
    if left is None:
        left = (SLIDE_W - Inches(w_in)) / 2
    top = BODY_TOP - Inches(0.15) if top is None else top
    slide.shapes.add_picture(str(path), left, top, Inches(w_in), Inches(h_in))
    if caption:
        tb = slide.shapes.add_textbox(MARGIN, top + Inches(h_in + 0.05),
                                      SLIDE_W - 2 * MARGIN, Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = MUTED
    return Inches(h_in)


def add_table(slide, header, rows, col_widths_in, top=BODY_TOP, size=10.5,
              left=MARGIN):
    n_rows, n_cols = len(rows) + 1, len(header)
    width = Inches(sum(col_widths_in))
    height = Inches(0.34 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = shape.table
    for j, w in enumerate(col_widths_in):
        t.columns[j].width = Inches(w)
    for j, h in enumerate(header):
        c = t.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = HDR_BG
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
                r.font.bold = True
                r.font.color.rgb = INK
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.color.rgb = INK
                    if j == 0:
                        r.font.name = "Consolas"
    return shape


# ---------------------------------------------------------------- 1 · title
s = new_slide()
tb = s.shapes.add_textbox(MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(2.2))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Roman grism line-grid simulation package"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = ACCENT
p2 = tb.text_frame.add_paragraph()
p2.text = "v0.13.0 rerun addendum — waves 1b/2/3 regenerated after the " \
          "sky→FPA placement fixes, with before/after comparison"
for r in p2.runs:
    r.font.size = Pt(16)
    r.font.color.rgb = INK
add_text(s, [
    "ADDENDUM to line_grid_wave1b_draft.pptx · 2026-07-31",
    "roman_disperser · branch feature/optical-model-line-test · commit 3d94b0f "
    "(merges main 12b98fc = v0.13.0)",
    "Stores: /mnt/roman-science/grs/line-tests-20260731{,-cont,-gal} · "
    "identical catalogs to 20260724 · full details in briefing_v0130_rerun.md",
], top=Inches(5.5), size=13, color=MUTED)

# ------------------------------------------------ 2 · what changed and why
s = new_slide("What was rerun, and why")
add_text(s, [
    "Three code fixes landed after the 20260724 campaign, all in the sky→FPA "
    "source-placement path:",
    ("v0.11.0 — precision: TF32 rotation on GPU and a float32 RA downcast", 1),
    ("v0.12.0 — projection: flat-sky approximation replaced by gnomonic TAN", 1),
    ("v0.13.0 — reproducibility: per-SCA RNG keys, CODEVER/GITSHA provenance "
     "in every product", 1),
    "All three waves were regenerated with bit-identical catalogs (verified "
    "before submission), so every difference is attributable to the code.",
    "36 SLURM jobs (7011–7046): 12 sims (a10g) + 12 romanisim L2 wraps "
    "(CRDS_CONTEXT = roman_0058.pmap, as in the originals) + 12 closures. "
    "No failures.",
    "Truth tables now self-identify: codever 0.13.0 + git SHA embedded in the "
    "parquet metadata (roman_disperser_provenance) and provenance JSON.",
], size=15, line_gap=9)

# ------------------------------------------------------- 3 · blast radius
s = new_slide("Reported science metrics barely move")
add_figure(s, ANALYSIS / "fig3_science_stability.png", max_h_in=4.9,
           max_w_in=5.6, left=MARGIN)
add_text(s, [
    "Per-SCA median d_disp, 20260724 vs 20260731: every point on the "
    "identity line.",
    "Max |Δ| over 18 SCAs:",
    ("stpsf pa000: 0.0059 px · stpsf pa010: 0.0046 px", 1),
    ("gauss pa000: 0.0047 px · gauss pa010: 0.0029 px", 1),
    "Reading: the package was internally self-consistent — simulation and "
    "prediction shared the placement code, so its errors cancelled in "
    "d_disp. The shipped wave conclusions were robust.",
    "The ~0.005 px residual deltas are edge-source reshuffling (row counts "
    "25,680 → 25,695 at PA 0) plus estimator noise.",
], left=Inches(6.5), top=Inches(1.5), width=Inches(6.3), size=14, line_gap=8)

# ------------------------------------------------------- 4 · placement map
s = new_slide("Absolute placement moved up to ~7 px at the field edge")
add_figure(s, ANALYSIS / "fig1_shift_map.png", max_h_in=5.5, max_w_in=5.9,
           left=MARGIN)
add_figure(s, ANALYSIS / "fig2_shift_cdf.png", max_h_in=2.9, max_w_in=6.4,
           left=Inches(6.6), top=Inches(1.3))
add_text(s, [
    "|shift|: median 1.82 px, p95 4.87 px, max 6.82 px "
    "(per-SCA medians 0.54–3.21 px; largest on SCA 18).",
    "Grows with field radius — the same radial signature the independent "
    "validation measured on the defect.",
    "Matches the predicted blast radius (1.84 px median, up to 7.1 px).",
    "Consumers dispersing from the shipped absolute positions should move "
    "to the 20260731 stores.",
], left=Inches(6.6), top=Inches(4.45), width=Inches(6.2), size=13, line_gap=7)

# ------------------------------------------------------ 5 · smoothness gate
s = new_slide("Independent-validation smoothness defect: eliminated")
h = add_figure(s, ANALYSIS / "fig4_smoothness.png", max_h_in=4.4,
               caption="Per-SCA rms of the residual after a degree-3 "
                       "polynomial fit of (RA, Dec) → (xsca, ysca), wave-1b "
                       "stpsf PA 0 truth tables. Pre-fix: 0.39–1.09 px "
                       "(the audit's non-smooth map). v0.13.0: "
                       "1.1e-4–3.6e-4 px — smooth at float precision.")
add_text(s, [
    "The audit's one real defect (0.6–2.9 px rms non-smoothness, surviving "
    "affine / degree-6 / TAN+SIP fits) was exactly the v0.11/v0.12 placement "
    "bugs. Its open question — which code places sources — is answered: "
    "get_fpa_pos.",
], top=Inches(6.55), size=13)

# ------------------------------------------ 6 · replication + provenance
s = new_slide("Waves 2 and 3 replicate")
add_text(s, ["Same estimators and comparisons as the 20260724 campaign, "
             "on the regenerated stores:"], size=14)
add_table(s, ["check", "20260724", "20260731", "reading"],
          [["wave-2 continuum null: max per-SCA |Δ median| vs wave 1b",
            "≤ 0.0046 px", "≤ 0.0045 px",
            "continuum still does not move line centroids"],
           ["wave-3 star→galaxy per-SCA correlation (STPSF)",
            "r = 0.998", "r = 0.998 / 0.996",
            "per-SCA pattern transfers to Sersic galaxies"],
           ["wave-3 max |galaxy − star| per-SCA median",
            "≤ 0.0064 px", "≤ 0.0080 px", "same order; absolute agreement"],
           ["wave-3 Gaussian control", "r ≈ 0.07–0.18", "r ≈ 0.08–0.09",
            "no per-SCA pattern exists to transfer (medians ~0.001 px "
            "noise); r not meaningful, unchanged"]],
          [4.6, 1.7, 1.8, 4.1], top=Inches(1.75), size=11)
add_text(s, [
    "Not compared: ISIM noise realisations — v0.13.0's per-SCA keys change "
    "every draw by construction, and GPU scatter-add is non-deterministic "
    "(roman_disperser#22). All gates are MODEL images with tolerances.",
    "Analysis: workbench/20260731-ssc-line-grid-analysis/before_after.py; "
    "figures, per-SCA tables, summary.json and this deck live in "
    "line-tests-20260731/analysis/.",
], top=Inches(4.6), size=13, line_gap=8)

prs.save(OUT)
print(f"wrote {OUT}")
