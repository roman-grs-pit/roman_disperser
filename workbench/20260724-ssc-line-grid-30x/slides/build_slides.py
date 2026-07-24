"""Build the wave-1b line-grid hand-off deck (draft).

Generates line_grid_wave1b_draft.pptx from the figures in figs/ (copies of the
2026-07-24 wave-1b briefing figures). Content mirrors the wave-1b briefing
(research log, 2026-07-24) with descriptive text only.

Not a project dependency: python-pptx runs from an ephemeral pixi env,

    pixi exec --spec python=3.12 --spec python-pptx python build_slides.py

Every number here is quoted from the wave-1b briefing / its figure scripts;
regenerate those first if the run store changes.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from PIL import Image

HERE = Path(__file__).parent
FIGS = HERE / "figs"
OUT = HERE / "line_grid_wave1b_draft.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
BODY_TOP = Inches(1.25)
INK = RGBColor(0x20, 0x20, 0x20)
MUTED = RGBColor(0x59, 0x59, 0x59)
ACCENT = RGBColor(0x00, 0x50, 0x80)
HDR_BG = RGBColor(0xE8, 0xEE, 0xF4)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide(title=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(MARGIN, Inches(0.35), SLIDE_W - 2 * MARGIN,
                                  Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    return s


def add_text(slide, lines, left=MARGIN, top=BODY_TOP, width=None,
             size=14, mono=False, color=INK, line_gap=6):
    """lines: list of (text, level) or plain strings (level 0)."""
    width = width or (SLIDE_W - 2 * MARGIN)
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = lvl
        p.space_after = Pt(line_gap)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            if mono:
                r.font.name = "Consolas"
    return tb


def add_figure(slide, path, caption=None, max_h_in=5.35):
    img = Image.open(path)
    w_px, h_px = img.size
    max_w_in = 12.2
    scale = min(max_w_in / (w_px / 150), max_h_in / (h_px / 150))
    w_in, h_in = (w_px / 150) * scale, (h_px / 150) * scale
    left = (SLIDE_W - Inches(w_in)) / 2
    slide.shapes.add_picture(str(path), left, BODY_TOP - Inches(0.15),
                             Inches(w_in), Inches(h_in))
    if caption:
        tb = slide.shapes.add_textbox(MARGIN, BODY_TOP + Inches(h_in - 0.05),
                                      SLIDE_W - 2 * MARGIN, Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = MUTED


def add_table(slide, header, rows, col_widths_in, top=BODY_TOP, size=10.5,
              left=MARGIN):
    n_rows, n_cols = len(rows) + 1, len(header)
    width = Inches(sum(col_widths_in))
    height = Inches(0.34 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = shape.table
    for j, w in enumerate(col_widths_in):
        t.columns[j].width = Inches(w)
    for j, h in enumerate(header):
        c = t.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = HDR_BG
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
                r.font.bold = True
                r.font.color.rgb = INK
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.color.rgb = INK
                    if j == 0:
                        r.font.name = "Consolas"
    return shape


# ---------------------------------------------------------------- 1 · title
s = new_slide()
tb = s.shapes.add_textbox(MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(2.2))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Roman grism line-grid simulation package"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = ACCENT
p2 = tb.text_frame.add_paragraph()
p2.text = "Wave 1b — emission-line point sources with exact truth tables, " \
          "for line-by-line validation of grism dispersion and extraction"
for r in p2.runs:
    r.font.size = Pt(16)
    r.font.color.rgb = INK
add_text(s, [
    "DRAFT · 2026-07-24",
    "roman_disperser · branch feature/optical-model-line-test · commit eee413e",
], top=Inches(5.7), size=13, color=MUTED)

# ------------------------------------------------------- 2 · file locations
s = new_slide("Where the data is")
add_text(s, [
    "Run store (Roman AWS cluster; readable from head and compute nodes):",
], size=14)
add_text(s, [
    "/mnt/roman-science/grs/line-tests-20260724/",
    "├── catalog_lines/     408K   input catalog: positions + line SEDs + line list",
    "├── lines_stpsf/       4.5G   simulations, real (STPSF) PSFs — one dir per PA",
    "├── lines_gauss/       4.5G   simulations, Gaussian-control PSFs — same layout",
    "├── lines_stpsf_l2/    7.3G   romanisim L2 ASDF products for lines_stpsf",
    "├── lines_gauss_l2/    7.3G   same for lines_gauss",
    "├── scripts/            36K   verbatim campaign drivers + configs + pointings",
    "├── slurm-meta/        148K   per-job audit: task script + environment record",
    "└── logs/              2.5M   per-file romanisim wrap logs",
], top=Inches(1.7), size=13, mono=True, line_gap=3)
add_text(s, [
    "Pointing directories are named pointing_{pa000,pa010}_001.001.001.001.001.001 "
    "(suffix = APT plan.pass.segment.observation.visit.exposure).",
    "Code: roman_disperser repo, branch feature/optical-model-line-test (PR #17).",
], top=Inches(5.6), size=13)

# ------------------------------------------------- 3 · sources and spectra
s = new_slide("What was simulated — sources and spectra")
add_text(s, [
    "15,112 point sources: every star with F158 ≤ 20 in the parent reference catalog.",
    "All stars carry the same spectrum: five Gaussian emission lines on a flat "
    "continuum that is then removed — the images contain lines only.",
    "Continuum anchor: mag 16.3072 in F158 (flux_scale = 3.0000e-7 maggies); "
    "line amplitudes are defined relative to that (removed) continuum.",
], size=14)
add_table(s, ["line_id", "λ [Å, vacuum]", "FWHM [Å]", "amp / continuum",
              "line flux [erg s⁻¹ cm⁻²]", "detected e⁻ (+1 order, 190.22 s)"],
          [["0", "11000", "10", "5.00", "1.44e-14", "22,600"],
           ["1", "12500", "10", "6.25", "1.39e-14", "32,000"],
           ["2", "15500", "10", "7.50", "1.09e-14", "27,500"],
           ["3", "17000", "10", "8.75", "1.05e-14", "24,800"],
           ["4", "19000", "10", "10.00", "9.63e-15", "19,300"]],
          [1.1, 1.7, 1.2, 1.8, 3.0, 3.4], top=Inches(3.15), size=12)
add_text(s, [
    "Detected e⁻ = 33×33 px aperture sum on the noise-free image × exposure time.",
], top=Inches(5.6), size=12, color=MUTED)

# ---------------------------------------------------- 4 · run matrix
s = new_slide("What was simulated — run matrix and observation")
add_text(s, [
    "Four runs = {STPSF PSFs, Gaussian PSFs} × {PA 0°, PA 10°}, all 18 SCAs each:",
    ("STPSF: real chromatic, field-dependent, asymmetric PSFs "
     "(4×4 spatial grid × 56 wavelengths per SCA, 4× oversampled)", 1),
    ("Gaussian: centred symmetric Gaussians (FWHM 2.5 px) in the same cache "
     "format — a PSF-shape control with identical geometry", 1),
    ("PA 10°: the same pointing rolled by 10° — a sky-orientation control", 1),
    "One pointing: RA 10.0°, Dec 0.0°, exposure 190.22 s, MA table 1036 "
    "(GRISM, 190.22287 s).",
    "Spectral orders 0, +1, +2 are dispersed for every source.",
    "All randomness is seeded (config seed 42); per-SCA RNG keys are recorded "
    "in the FITS headers and meta files.",
], size=15, line_gap=9)

# ---------------------------------------------------- 5 · pipeline
s = new_slide("What was simulated — pipeline")
add_text(s, [
    "1.  Simulate.  JAX-based disperser renders all sources, orders and SCAs.",
    ("Output per SCA: FITS with a noise-free MODEL extension (counts/s) and a "
     "Poisson-sampled ISIM extension (counts at 190.22 s).", 1),
    "2.  Wrap.  romanisim converts each MODEL/ISIM to a Roman L2 ASDF "
    "(ramp simulation with MA table 1036, then ramp fit).",
    ("Output per SCA: *_l2.asdf with data / err / dq (DN/s), 11 resultants, "
     "effective exposure 190.22287 s.", 1),
    "3.  Measure + truth.  Every first-order (star, line) position is predicted "
    "by the optical model and measured on the noise-free MODEL image; "
    "independently, float64 truth positions are tabulated for all orders.",
    ("Output per pointing: residuals.parquet (closure measurement) and "
     "truth_lines.parquet (shipped truth table).", 1),
], size=15, line_gap=9)

# ---------------------------------------------------- 6-7 · figures sim/L2
s = new_slide("Simulated image (SCA 1, STPSF, PA 0)")
add_figure(s, FIGS / "fig1-simulation.png", max_h_in=5.1,
           caption="Noise-free MODEL image. Top: full SCA (8×8 max-pooled for "
                   "display). Bottom: one star's +1-order trace, transposed to "
                   "equal aspect — five isolated emission-line spots with their "
                   "17×17 px measurement boxes.")

s = new_slide("L2 product and noise (same strip)")
add_figure(s, FIGS / "fig2-l2-noise.png", max_h_in=5.1,
           caption="L2 data for the same strip (background 1.96 DN/s, robust "
                   "σ 0.105 DN/s). Per-line peak SNR 54–84, box SNR 73–110; "
                   "median peak SNR over all 322 measured boxes on SCA 1 is 66.")

# ---------------------------------------------------- 8 · measurement
s = new_slide("Line-position measurement")
add_text(s, [
    "For every first-order (star, line) pair, on the noise-free MODEL image:",
    ("Predict the line's detector position with the optical model "
     "(two independent implementations; they agree to ≤ 0.005 px).", 1),
    ("Cut a 17×17 px box at the prediction; drop the pair if the box is not "
     "fully on-detector (this is the only cut).", 1),
    ("Collapse a ±3 px band across dispersion; fit and subtract a linear "
     "baseline to the off-line pixels (|Δ| > 4 px).", 1),
    ("Measure the flux-weighted centroid within ±4 px of the prediction.", 1),
    "Residuals d_disp / d_cross are measured − predicted, projected on the "
    "local dispersion / cross-dispersion directions, in native pixels.",
    "Convention: a “measured position” is a flux-weighted centroid. For the "
    "STPSF PSFs this differs from the PSF peak by up to ~0.08 px, and depends "
    "on the window (±4 px is core-weighted; wing-weighted apertures move "
    "per-SCA offsets by up to ~30%).",
], size=15, line_gap=8)

s = new_slide("The estimator on one line")
add_figure(s, FIGS / "fig3-estimator.png", max_h_in=5.1,
           caption="One 17×17 box (1.55 µm line, SCA 1): the ±3 px collapse "
                   "band, the band-summed profile, linear baseline, ±4 px "
                   "centroid window, and measured centroid vs prediction.")

# ---------------------------------------------------- 10 · results table
s = new_slide("Results — summary")
add_text(s, ["Statistics over all finite-centroid first-order boxes "
             "(MAD = unscaled median absolute deviation):"], size=14)
add_table(s, ["run", "boxes (measured/total)", "d_disp median [px]",
              "d_disp MAD [px]", "d_cross median [px]", "optical-model "
              "implementations agree to [px]"],
          [["stpsf  pa000", "5290 / 8575", "−0.015", "0.041", "+0.024", "0.0040"],
           ["stpsf  pa010", "5266 / 8465", "−0.017", "0.040", "+0.025", "0.0053"],
           ["gauss  pa000", "5290 / 8575", "+0.000", "0.010", "−0.001", "0.0040"],
           ["gauss  pa010", "5266 / 8465", "−0.000", "0.010", "−0.001", "0.0053"]],
          [1.7, 2.3, 2.0, 1.8, 2.0, 2.4], top=Inches(1.9), size=12)
add_text(s, [
    "Unmeasured boxes are all geometric edge drops (17×17 box not fully "
    "on-detector); none are dropped for lack of flux.",
    "Every per-SCA median agrees with the 1× wave-1 run to < 1e-5 px, and with "
    "the 2026-07-03 validation run to ≤ 0.0036 px.",
], top=Inches(4.1), size=13)

# ---------------------------------------------------- 11-14 · result figures
s = new_slide("Results — per-SCA median offset")
add_figure(s, FIGS / "fig4-per-sca-medians.png", max_h_in=5.0,
           caption="Median d_disp per SCA for the four runs, with the "
                   "2026-07-03 reference overlaid. STPSF runs show a fixed "
                   "per-SCA pattern up to ±0.065 px, the same at both PAs; "
                   "Gaussian-control runs are consistent with zero on every SCA.")

s = new_slide("Results — per-SCA distributions")
add_figure(s, FIGS / "fig5-per-sca-violins.png", max_h_in=5.0,
           caption="d_disp distributions per SCA (PA 0 runs; PA 10 medians as "
                   "open triangles). Per-SCA MADs: 0.011–0.028 px (STPSF), "
                   "0.008–0.014 px (Gaussian).")

s = new_slide("Results — offset vs wavelength")
add_figure(s, FIGS / "fig6-line-trend.png", max_h_in=5.0,
           caption="Median d_disp per line, pooled over SCAs (opposite-sign "
                   "SCAs partially cancel). STPSF runs trend from −0.008 px at "
                   "1.10 µm to −0.021 px at 1.90 µm; Gaussian runs are flat "
                   "within ±0.003 px.")

s = new_slide("Results — comparison with PSF flux centroids")
add_figure(s, FIGS / "fig8-psf-centroids.png", max_h_in=5.0,
           caption="Measured per-(SCA, line) median offsets vs flux-weighted "
                   "centroids computed directly from the PSF cache stamps with "
                   "the same band/baseline/window procedure: r = 0.973, "
                   "slope 0.97, rms 0.011 px. The Gaussian-cache centroids are "
                   "zero to machine precision.")

# ---------------------------------------------------- 15 · file guide intro
s = new_slide("File guide — per-pointing contents")
add_text(s, ["Inside each pointing directory (lines_{stpsf,gauss}/pointing_*/):"],
         size=14)
add_text(s, [
    "grism_*_detSCA{NN}.fits       per-SCA image: MODEL (counts/s) + ISIM (counts)",
    "grism_*_detSCA{NN}.png        quicklook PNG per SCA (+ one focal-plane mosaic)",
    "grism_*_sources.parquet       manifest of every (source, SCA, order) rendered",
    "grism_*_meta.yaml             pointing, exposure, RNG keys, batching, git SHA",
    "residuals.parquet             measured vs predicted line positions (+1 order)",
    "truth_lines.parquet           float64 truth positions, all orders and lines",
    "truth_lines_provenance.json   writer script, commit, inputs, row counts",
], top=Inches(1.7), size=13, mono=True, line_gap=4)
add_text(s, [
    "L2 trees (lines_*_l2/pointing_*/) hold one *_detSCA{NN}_l2.asdf per SCA.",
    "catalog_lines/ holds the shared input catalog (next slides).",
    "The following slides give the columns of each tabular file.",
], top=Inches(4.75), size=13)

# ---------------------------------------------------- 16 · metadata.parquet
s = new_slide("File: catalog_lines/metadata.parquet")
add_text(s, ["One row per source (15,112 rows) — the simulation input catalog."],
         size=13)
add_table(s, ["column", "type", "units", "meaning"],
          [["ra, dec", "float64", "deg", "source position (ICRS)"],
           ["type", "str", "—", "\"PSF\" for all rows (point sources)"],
           ["n, half_light_radius, pa, ba", "float32", "mixed",
            "Sérsic morphology; trivial (0/0/0/1) for stars"],
           ["F158", "float32", "maggies",
            "apparent F158 flux, linear AB units; 3.0e-7 (mag 16.31) for every star"],
           ["z_obs, z_cosmo", "float32", "—", "redshifts; 0 for stars"],
           ["sed_index", "int32", "—", "row index into seds.zarr/star_seds"],
           ["flux_scale", "float32", "maggies",
            "SED multiplier applied at dispersal; equals F158 for stars"],
           ["sim", "int16", "—", "SED-storage partition; 0 for stars"],
           ["src_index", "int32", "—",
            "row in the parent reference catalog — cross-run provenance key"]],
          [2.6, 0.9, 0.9, 7.8], top=Inches(1.75), size=10.5)

# ---------------------------------------------------- 17 · seds + lines
s = new_slide("Files: catalog_lines/seds.zarr and lines.ecsv")
add_text(s, [
    "seds.zarr — the line-only SED, f_λ (FLAM) on a vacuum-wavelength grid "
    "9000–21000 Å in 2 Å steps (6001 samples); one shared template, "
    "referenced by sed_index. Scale by flux_scale for physical flux.",
    "lines.ecsv — the injected line list:",
], size=13, line_gap=6)
add_table(s, ["column", "units", "meaning"],
          [["line_id", "—",
            "0–4, blue to red; key used by residuals and truth tables"],
           ["center_A", "Å (vacuum)", "11000 / 12500 / 15500 / 17000 / 19000"],
           ["fwhm_A", "Å", "10.0 for all lines"],
           ["amp_rel", "—",
            "line peak / local continuum at SED level: 5.0–10.0 "
            "(the continuum itself was removed; no_continuum = true)"]],
          [1.6, 1.6, 9.0], top=Inches(2.7), size=11)
add_text(s, ["provenance.json records the builder inputs: parent catalog, "
             "mag limit, continuum mag, line parameters, wavelength grid, "
             "build timestamp."], top=Inches(4.8), size=12, color=MUTED)

# ---------------------------------------------------- 18 · FITS
s = new_slide("File: grism_*_detSCA{NN}.fits")
add_table(s, ["extension", "content"],
          [["0  PRIMARY", "header only (metadata below)"],
           ["1  MODEL", "noise-free count-rate image, counts/s, float32, 4088×4088"],
           ["2  ISIM", "Poisson-sampled counts at 190.22 s, float32, 4088×4088"]],
          [2.2, 9.0], top=BODY_TOP, size=11.5)
add_text(s, ["Primary-header keywords:"], top=Inches(2.75), size=13)
add_table(s, ["keyword", "meaning"],
          [["WFICENRA / WFICENDEC / WFICENPA", "pointing RA, Dec, PA [deg]"],
           ["DETNUM", "SCA number 1–18"],
           ["EXPTIME", "exposure time [s]"],
           ["MA_TABLE", "MA table number (1036)"],
           ["SEED / RNDSEED0 / RNDSEED1", "top seed + per-SCA JAX RNG key words"],
           ["GITSHA", "pipeline git commit"],
           ["PLAN, PASS, SEGMENT, OBS, VISIT, EXPOSURE", "APT identifiers"]],
          [3.6, 7.6], top=Inches(3.15), size=10.5)

# ---------------------------------------------------- 19 · sources.parquet
s = new_slide("File: grism_*_sources.parquet")
add_text(s, ["The dispersal manifest: one row per (source, SCA, order) "
             "actually rendered."], size=13)
add_table(s, ["column", "type", "meaning"],
          [["catalog_index", "int64",
            "row index into this run's catalog_lines/metadata.parquet"],
           ["sca", "int64", "SCA number 1–18"],
           ["order", "str", "spectral order: \"0\", \"1\", or \"2\""],
           ["type", "str", "\"PSF\" (all rows)"],
           ["xsca, ysca", "float64",
            "undispersed source position on this SCA [1-indexed FITS px]"],
           ["ra, dec", "float64", "source position [deg]"],
           ["flux_scale", "float64", "SED multiplier applied at dispersal"],
           ["F158", "float64", "F158 flux [maggies]; same on every row of a source"]],
          [2.0, 1.0, 9.2], top=Inches(1.85), size=10.5)

# ---------------------------------------------------- 20 · residuals.parquet
s = new_slide("File: residuals.parquet")
add_text(s, ["The closure measurement: one row per first-order (star, line) "
             "box — 8575 (PA 0) / 8465 (PA 10) rows."], size=13)
add_table(s, ["column", "type", "meaning"],
          [["sca", "int64", "SCA number"],
           ["catalog_index", "int64", "source key (into catalog metadata)"],
           ["xsca, ysca", "float64", "undispersed position [1-indexed FITS px]"],
           ["line_id, center_A", "int64/float64", "which line (see lines.ecsv)"],
           ["x_pred_jax, y_pred_jax", "float32",
            "JAX optical-model prediction [px]; float32 path, ULP ~1e-3 px"],
           ["x_pred_class, y_pred_class", "float64",
            "reference-implementation prediction [px]"],
           ["x_meas, y_meas", "float64",
            "measured centroid [px]; NaN if box off-detector"],
           ["dx, dy", "float32", "x_meas − x_pred_jax, y_meas − y_pred_jax [px]"],
           ["d_disp, d_cross", "float64",
            "residual on local dispersion / cross-dispersion axes [px] — "
            "the science columns"]],
          [2.4, 1.3, 8.5], top=Inches(1.85), size=10.5)

# ---------------------------------------------------- 21 · truth_lines
s = new_slide("File: truth_lines.parquet")
add_text(s, ["The shipped truth table: one row per (source, SCA, order, line), "
             "orders 0/1/2 — 25,680 rows; positions evaluated in float64."],
         size=13)
add_table(s, ["column", "type", "meaning"],
          [["src_index", "int32",
            "row in the parent reference catalog — cross-run provenance key"],
           ["catalog_index", "int64", "row in this run's catalog metadata"],
           ["ra, dec", "float64", "source position [deg]"],
           ["sca", "int16", "SCA number"],
           ["order", "str", "\"0\", \"1\", \"2\""],
           ["xsca, ysca", "float64", "undispersed position [1-indexed FITS px]"],
           ["line_id, center_A", "int64/float64", "which line [Å, vacuum]"],
           ["fwhm_A, amp_rel", "float64", "line width and relative amplitude"],
           ["x_pred, y_pred", "float64",
            "predicted dispersed line position [1-indexed FITS px]"],
           ["on_detector", "bool",
            "prediction point inside [0.5, 4088.5]² (weaker than the "
            "residuals' box-fully-on-detector cut)"]],
          [2.2, 1.3, 8.7], top=Inches(1.95), size=10.5)

# ---------------------------------------------------- 22 · L2 + audit
s = new_slide("Files: L2 ASDF, meta YAML, audit records")
add_text(s, [
    "*_detSCA{NN}_l2.asdf — romanisim L2 product per SCA (202 MB): "
    "roman.data / err / dq and var arrays (DN/s), full Roman metadata tree; "
    "MA table 1036, 11 resultants, effective exposure 190.22287 s.",
    "grism_*_meta.yaml — pointing, exposure, seed and per-SCA JAX RNG keys "
    "(reconstructable via jax.random.wrap_key_data), wavelength step, batching, "
    "per-SCA source counts by order, pipeline git SHA, APT identifiers.",
    "scripts/ — verbatim copy, at submission time, of the campaign driver, "
    "configs, pointing tables and truth-table writer.",
    "slurm-meta/ — per-job audit: the exact task script plus an .env record "
    "(job id, stage, partition, resources, config paths, git commits, "
    "timestamp).",
    "logs/ — per-file romanisim wrap logs.",
], size=14, line_gap=10)

prs.save(OUT)
print(f"wrote {OUT} ({len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides)")
